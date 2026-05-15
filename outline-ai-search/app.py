import asyncio
import hashlib
import os
import re
import sqlite3
import uuid
from contextlib import suppress
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from zipfile import BadZipFile

import httpx
import psycopg
from docx import Document
from docx.opc.exceptions import PackageNotFoundError
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from fastapi.responses import FileResponse, HTMLResponse
from pptx import Presentation
from pydantic import BaseModel
from pypdf import PdfReader
from qdrant_client import QdrantClient
from qdrant_client.http.models import (
    Distance,
    FieldCondition,
    Filter,
    MatchValue,
    PointStruct,
    VectorParams,
)

DATABASE_URL = os.getenv("DATABASE_URL", "postgres://user:pass@postgres:5432/outline")
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://ollama:11434")
QDRANT_URL = os.getenv("QDRANT_URL", "http://qdrant:6333")
EMBED_MODEL = os.getenv("EMBED_MODEL", "nomic-embed-text")
CHAT_MODEL = os.getenv("CHAT_MODEL", "gemma4:e2b")
CHAT_KEEP_ALIVE = -1
COLLECTION = os.getenv("QDRANT_COLLECTION", "outline_documents")
OUTLINE_URL = os.getenv("OUTLINE_URL", "http://localhost:8088")
AUTO_INDEX_INTERVAL_SECONDS = int(os.getenv("AUTO_INDEX_INTERVAL_SECONDS", "60"))
FAST_INDEX_INTERVAL_SECONDS = int(os.getenv("FAST_INDEX_INTERVAL_SECONDS", "30"))
DATA_DIR = Path(os.getenv("AI_DATA_DIR", "/data"))
OUTLINE_FILES_DIR = Path(os.getenv("OUTLINE_FILES_DIR", "/outline-data"))
FILES_DIR = DATA_DIR / "files"
DB_PATH = DATA_DIR / "ai_search.sqlite3"
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}
LEGACY_EXTENSIONS = {".doc", ".xls", ".ppt"}
MEDIA_TYPES = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
MIN_OUTLINE_TEXT_CHARS = int(os.getenv("MIN_OUTLINE_TEXT_CHARS", "20"))
MIN_SEARCH_SCORE = float(os.getenv("MIN_SEARCH_SCORE", "0.50"))
TITLE_MATCH_BOOST = float(os.getenv("TITLE_MATCH_BOOST", "0.18"))

app = FastAPI(title="Outline AI Search")
qdrant = QdrantClient(url=QDRANT_URL)
index_lock = asyncio.Lock()
auto_index_task: asyncio.Task[None] | None = None


class SearchRequest(BaseModel):
    query: str
    limit: int = 5
    collectionId: str | None = None
    allowedDocumentIds: list[str] | None = None


class ChatHistoryMessage(BaseModel):
    role: str
    content: str


class ChatRequest(BaseModel):
    message: str
    limit: int = 5
    collectionId: str | None = None
    allowedDocumentIds: list[str] | None = None
    history: list[ChatHistoryMessage] = []


class ConvertResponse(BaseModel):
    title: str
    text: str


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def db() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db() -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    FILES_DIR.mkdir(parents=True, exist_ok=True)
    with db() as conn:
        conn.execute(
            """
            create table if not exists files (
              id text primary key,
              filename text not null,
              title text not null,
              content_type text,
              size integer not null,
              path text not null,
              status text not null,
              chunks integer not null default 0,
              error text,
              collection_id text,
              collection_name text,
              created_at text not null,
              indexed_at text
            )
            """
        )
        conn.execute(
            """
            create table if not exists index_jobs (
              id text primary key,
              source text not null,
              source_id text,
              status text not null,
              documents integer not null default 0,
              chunks integer not null default 0,
              error text,
              created_at text not null,
              finished_at text
            )
            """
        )
        for statement in [
            "alter table files add column collection_id text",
            "alter table files add column collection_name text",
        ]:
            with suppress(sqlite3.OperationalError):
                conn.execute(statement)
        conn.execute(
            """
            create table if not exists attachment_index (
              id text primary key,
              document_id text not null,
              filename text not null,
              content_type text,
              size integer not null,
              status text not null,
              chunks integer not null default 0,
              error text,
              updated_at text,
              indexed_at text
            )
            """
        )
        conn.execute(
            """
            create table if not exists document_index (
              id text primary key,
              title text not null,
              collection text,
              status text not null,
              chunks integer not null default 0,
              error text,
              updated_at text,
              indexed_at text
            )
            """
        )


@app.on_event("startup")
async def startup() -> None:
    global auto_index_task
    init_db()
    auto_index_task = asyncio.create_task(auto_index_loop())
    asyncio.create_task(warmup_chat_model())


@app.on_event("shutdown")
async def shutdown() -> None:
    if auto_index_task:
        auto_index_task.cancel()
        with suppress(asyncio.CancelledError):
            await auto_index_task


def clean_text(value: str | None) -> str:
    if not value:
        return ""
    return re.sub(r"\s+", " ", value).strip()


def repair_mojibake(value: str | None) -> str:
    text = clean_text(value)
    for _ in range(2):
        if not any(marker in text for marker in ("Ð", "Ñ", "Â", "â")):
            break
        try:
            fixed = text.encode("latin1").decode("utf-8")
        except UnicodeError:
            break
        if fixed == text:
            break
        text = clean_text(fixed)
    return text


_RU_TRANSLIT = str.maketrans({
    "а": "a", "б": "b", "в": "v", "г": "g", "д": "d", "е": "e", "ё": "e",
    "ж": "zh", "з": "z", "и": "i", "й": "y", "к": "k", "л": "l", "м": "m",
    "н": "n", "о": "o", "п": "p", "р": "r", "с": "s", "т": "t", "у": "u",
    "ф": "f", "х": "h", "ц": "ts", "ч": "ch", "ш": "sh", "щ": "sch", "ъ": "",
    "ы": "y", "ь": "", "э": "e", "ю": "yu", "я": "ya",
})


def outline_slug(title: str | None) -> str:
    slug = repair_mojibake(title).lower().translate(_RU_TRANSLIT)
    slug = re.sub(r"[^a-z0-9_\s-]+", "", slug)
    slug = re.sub(r"[\s-]+", "-", slug).strip("-")[:80]
    return slug or "doc"


def document_url(title: str | None, url_id: str | None) -> str:
    return f"{OUTLINE_URL}/doc/{outline_slug(title)}-{url_id}"


def attachment_url(attachment_id: str | None) -> str:
    if not attachment_id:
        return OUTLINE_URL
    return f"{OUTLINE_URL}/api/attachments.redirect?id={attachment_id}"



def lookup_document_url(document_id: str | None) -> str | None:
    if not document_id:
        return None
    try:
        with psycopg.connect(DATABASE_URL) as conn:
            with conn.cursor() as cur:
                cur.execute('select title, "urlId" from documents where id = %s and "deletedAt" is null', (document_id,))
                row = cur.fetchone()
        if not row:
            return None
        return document_url(row[0], row[1])
    except Exception:
        return None

def normalize_match_text(value: str | None) -> str:
    value = clean_text(value).lower()
    value = re.sub(r"\.(pdf|docx|xlsx|pptx)$", "", value, flags=re.IGNORECASE)
    return re.sub(r"[^\w]+", " ", value, flags=re.UNICODE).strip()


def title_match_score(query: str, payload: dict[str, Any]) -> float:
    query_text = normalize_match_text(query)
    if not query_text:
        return 0.0
    title_text = normalize_match_text(payload.get("title"))
    filename_text = normalize_match_text(payload.get("filename"))
    haystacks = [value for value in (title_text, filename_text) if value]
    if any(query_text == value or query_text in value or value in query_text for value in haystacks):
        return TITLE_MATCH_BOOST
    query_tokens = {token for token in query_text.split() if len(token) > 3}
    if not query_tokens:
        return 0.0
    for value in haystacks:
        value_tokens = set(value.split())
        overlap = len(query_tokens & value_tokens)
        if overlap:
            return min(TITLE_MATCH_BOOST, 0.06 * overlap)
    return 0.0


def chunk_text(text: str, max_chars: int = 1800, overlap: int = 250) -> list[str]:
    text = clean_text(text)
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = max(0, end - overlap)
    return chunks


def useful_text_length(text: str) -> int:
    return sum(1 for char in clean_text(text) if char.isalnum())


async def warmup_chat_model() -> None:
    await asyncio.sleep(2)
    try:
        async with httpx.AsyncClient(timeout=180) as client:
            await client.post(
                f"{OLLAMA_URL}/api/chat",
                json={
                    "model": CHAT_MODEL,
                    "messages": [{"role": "user", "content": "Warm up. Reply OK."}],
                    "stream": False,
                    "think": False,
                    "keep_alive": CHAT_KEEP_ALIVE,
                    "options": {"num_ctx": 2048, "num_predict": 8, "temperature": 0},
                },
            )
    except Exception as exc:
        print(f"AI chat warmup failed: {exc}", flush=True)


async def embed(text: str) -> list[float]:
    async with httpx.AsyncClient(timeout=120) as client:
        response = await client.post(
            f"{OLLAMA_URL}/api/embeddings",
            json={"model": EMBED_MODEL, "prompt": text},
        )
        if response.status_code >= 400:
            raise HTTPException(status_code=502, detail=f"Ollama error: {response.text}")
        return response.json()["embedding"]


async def ensure_collection(vector_size: int) -> None:
    if any(item.name == COLLECTION for item in qdrant.get_collections().collections):
        return
    qdrant.create_collection(
        collection_name=COLLECTION,
        vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
    )


def point_id(source_type: str, source_id: str, chunk_index: int) -> str:
    return hashlib.md5(f"{source_type}:{source_id}:{chunk_index}".encode()).hexdigest()


def delete_source_points(source_type: str, source_id: str) -> None:
    try:
        qdrant.delete(
            collection_name=COLLECTION,
            points_selector=Filter(
                must=[
                    FieldCondition(key="sourceType", match=MatchValue(value=source_type)),
                    FieldCondition(key="sourceId", match=MatchValue(value=source_id)),
                ]
            ),
        )
    except Exception:
        return


def delete_source_type_points(source_type: str) -> None:
    try:
        qdrant.delete(
            collection_name=COLLECTION,
            points_selector=Filter(
                must=[
                    FieldCondition(key="sourceType", match=MatchValue(value=source_type)),
                ]
            ),
        )
    except Exception:
        return


def fetch_documents() -> list[dict[str, Any]]:
    query = """
        select d.id::text, d."urlId", d.title, d.text, d."updatedAt"::text, c.name, c.id::text
        from documents d
        left join collections c on c.id = d."collectionId"
        where d."deletedAt" is null
          and d."archivedAt" is null
        order by d."updatedAt" desc
    """
    with psycopg.connect(DATABASE_URL) as conn:
        with conn.cursor() as cur:
            cur.execute(query)
            rows = cur.fetchall()
    return [
        {
            "id": row[0],
            "urlId": row[1],
            "title": row[2],
            "text": row[3],
            "updatedAt": row[4],
            "collection": row[5],
            "collectionId": row[6],
        }
        for row in rows
    ]


def fetch_indexable_attachments() -> list[dict[str, Any]]:
    query = """
        select a.id::text, a."documentId"::text, a.key, a."contentType", a.size,
               a."updatedAt"::text, d.title, d."urlId", c.name, c.id::text
        from attachments a
        join documents d on d.id = a."documentId"
        left join collections c on c.id = d."collectionId"
        where (
            a."contentType" in (
              'application/pdf',
              'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
              'application/vnd.openxmlformats-officedocument.presentationml.presentation',
              'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            )
            or lower(a.key) like '%.pdf'
            or lower(a.key) like '%.docx'
            or lower(a.key) like '%.pptx'
            or lower(a.key) like '%.xlsx'
        )
          and d."deletedAt" is null
          and d."archivedAt" is null
        order by a."updatedAt" desc
    """
    with psycopg.connect(DATABASE_URL) as conn:
        with conn.cursor() as cur:
            cur.execute(query)
            rows = cur.fetchall()
    return [
        {
            "id": row[0],
            "documentId": row[1],
            "key": row[2],
            "contentType": row[3],
            "size": row[4],
            "updatedAt": row[5],
            "documentTitle": row[6],
            "urlId": row[7],
            "collection": row[8],
            "collectionId": row[9],
        }
        for row in rows
    ]


def fetch_collections() -> list[dict[str, str]]:
    query = """
        select id::text, name, "urlId"
        from collections
        where "deletedAt" is null
        order by name asc
    """
    with psycopg.connect(DATABASE_URL) as conn:
        with conn.cursor() as cur:
            cur.execute(query)
            rows = cur.fetchall()
    return [{"id": row[0], "name": row[1], "urlId": row[2]} for row in rows]


def get_collection(collection_id: str | None) -> dict[str, str] | None:
    if not collection_id:
        return None
    query = """
        select id::text, name, "urlId"
        from collections
        where id = %s and "deletedAt" is null
        limit 1
    """
    with psycopg.connect(DATABASE_URL) as conn:
        with conn.cursor() as cur:
            cur.execute(query, (collection_id,))
            row = cur.fetchone()
    if not row:
        return None
    return {"id": row[0], "name": row[1], "urlId": row[2]}


def expand_query(query: str) -> str:
    normalized = query.lower()
    if "РєС†Рї" in normalized:
        return f"{query}\nРљРѕРјР°РЅРґР° С†РёС„СЂРѕРІРѕРіРѕ РїСЂРѕСЂС‹РІР° РєРѕРЅРєСѓСЂСЃ РїРѕР»РѕР¶РµРЅРёРµ"
    return query


def unique_sources(results: list[dict[str, Any]]) -> list[dict[str, Any]]:
    seen: set[Any] = set()
    sources: list[dict[str, Any]] = []
    for result in results:
        key = source_key(result)
        if key in seen:
            continue
        seen.add(key)
        sources.append(result)
    return sources


def readable_source(source: dict[str, Any]) -> dict[str, Any]:
    result = dict(source)
    if result.get("sourceType") == "outline_attachment":
        page_url = result.get("documentUrl") or lookup_document_url(result.get("documentId"))
        if page_url:
            result["url"] = page_url
        if result.get("filename"):
            result["title"] = result["filename"]
        elif isinstance(result.get("title"), str):
            result["title"] = re.split(r"\s+(?:?|??|????)\s+", result["title"])[-1]
    return result

def readable_sources(results: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [readable_source(source) for source in unique_sources(results)]


def readable_answer(answer: str) -> str:
    answer = answer.replace("**", "")
    answer = re.sub(r"https?://\S+", "", answer)
    answer = re.sub(r"\[РСЃС‚РѕС‡РЅРёРє:[^\]]+\]", "", answer, flags=re.IGNORECASE)
    answer = re.sub(r"\s+\n", "\n", answer)
    return clean_text(answer)


def context_results(results: list[dict[str, Any]], limit: int = 3) -> list[dict[str, Any]]:
    selected: list[dict[str, Any]] = []
    seen: set[Any] = set()
    for result in results:
        if result.get("score", 0) < MIN_SEARCH_SCORE and not result.get("titleMatch"):
            continue
        key = source_key(result)
        if key in seen:
            continue
        seen.add(key)
        selected.append(result)
        if len(selected) >= limit:
            break
    return selected


def focused_results(results: list[dict[str, Any]], limit: int = 3) -> list[dict[str, Any]]:
    if not results:
        return []
    first = results[0]
    document_id = first.get("documentId")
    if not document_id:
        return context_results(results, limit=limit)
    if first.get("score", 0) < MIN_SEARCH_SCORE and not first.get("titleMatch"):
        return []
    same_document = [item for item in results if item.get("documentId") == document_id]
    focused = context_results(same_document, limit=limit)
    return focused or context_results(results, limit=limit)


def source_key(source: dict[str, Any]) -> Any:
    title = clean_text(source.get("title"))
    title = re.sub(r"\.(pdf|docx|xlsx|pptx)$", "", title, flags=re.IGNORECASE)
    for separator in (" · ", " В· "):
        if separator in title:
            title = title.split(separator, 1)[0]
            break
    normalized = re.sub(r"[^\w]+", " ", title.lower(), flags=re.UNICODE).strip()
    if normalized and normalized not in {"1", "doc"} and len(normalized) > 3:
        return ("title", normalized)
    return (
        "id",
        source.get("documentId")
        or source.get("sourceId")
        or source.get("url")
        or source.get("title"),
    )


def extract_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    pages: list[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text())
        if text:
            pages.append(f"Р РЋРЎвЂљРЎР‚Р В°Р Р…Р С‘РЎвЂ Р В° {page_number}\n{text}")
    return "\n\n".join(pages)


def extract_docx_text(path: Path) -> str:
    document = Document(str(path))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = clean_text(paragraph.text)
        if text:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            values = [clean_text(cell.text) for cell in row.cells]
            values = [value for value in values if value]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def extract_xlsx_text(path: Path) -> str:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"Р вЂєР С‘РЎРѓРЎвЂљ: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [clean_text(str(value)) for value in row if value is not None and clean_text(str(value))]
            if values:
                parts.append(" | ".join(values))
    workbook.close()
    return "\n".join(parts)


def extract_pptx_text(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = clean_text(shape.text)
            if text:
                slide_parts.append(text)
        if slide_parts:
            parts.append(f"Slide {slide_number}\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def extract_file_text(path: Path) -> str:
    extension = path.suffix.lower()
    try:
        if extension == ".pdf":
            return extract_pdf_text(path)
        if extension == ".docx":
            return extract_docx_text(path)
        if extension == ".xlsx":
            return extract_xlsx_text(path)
        if extension == ".pptx":
            return extract_pptx_text(path)
    except (BadZipFile, InvalidFileException, PackageNotFoundError, KeyError) as exc:
        raise ValueError(
            "Р В¤Р В°Р в„–Р В» Р Р†РЎвЂ№Р С–Р В»РЎРЏР Т‘Р С‘РЎвЂљ Р С—Р С•Р Р†РЎР‚Р ВµР В¶Р Т‘Р ВµР Р…Р Р…РЎвЂ№Р С Р С‘Р В»Р С‘ Р Р…Р Вµ РЎРѓР С•Р С•РЎвЂљР Р†Р ВµРЎвЂљРЎРѓРЎвЂљР Р†РЎС“Р ВµРЎвЂљ РЎРѓР Р†Р С•Р ВµР СРЎС“ РЎР‚Р В°РЎРѓРЎв‚¬Р С‘РЎР‚Р ВµР Р…Р С‘РЎР‹. "
            "Р вЂўРЎРѓР В»Р С‘ РЎРЊРЎвЂљР С• РЎРѓРЎвЂљР В°РЎР‚РЎвЂ№Р в„– .doc/.xls, РЎРѓР С•РЎвЂ¦РЎР‚Р В°Р Р…Р С‘РЎвЂљР Вµ Р ВµР С–Р С• Р С”Р В°Р С” .docx/.xlsx Р С‘ Р В·Р В°Р С–РЎР‚РЎС“Р В·Р С‘РЎвЂљР Вµ Р В·Р В°Р Р…Р С•Р Р†Р С•."
        ) from exc
    raise ValueError("Р СџР С•Р Т‘Р Т‘Р ВµРЎР‚Р В¶Р С‘Р Р†Р В°РЎР‹РЎвЂљРЎРѓРЎРЏ РЎвЂљР С•Р В»РЎРЉР С”Р С• PDF, DOCX Р С‘ XLSX.")


async def index_chunks(
    *,
    source_type: str,
    source_id: str,
    title: str,
    url: str,
    text: str,
    extra_payload: dict[str, Any] | None = None,
) -> int:
    chunks = chunk_text(text)
    delete_source_points(source_type, source_id)
    indexed = 0
    for index, chunk in enumerate(chunks):
        vector = await embed(chunk)
        await ensure_collection(len(vector))
        payload = {
            "sourceType": source_type,
            "sourceId": source_id,
            "title": title,
            "url": url,
            "chunk": chunk,
            "chunkIndex": index,
        }
        if extra_payload:
            payload.update(extra_payload)
        qdrant.upsert(
            collection_name=COLLECTION,
            points=[
                PointStruct(
                    id=point_id(source_type, source_id, index),
                    vector=vector,
                    payload=payload,
                )
            ],
        )
        indexed += 1
    return indexed


def create_job(source: str, source_id: str | None = None) -> str:
    job_id = str(uuid.uuid4())
    with db() as conn:
        conn.execute(
            "insert into index_jobs (id, source, source_id, status, created_at) values (?, ?, ?, 'running', ?)",
            (job_id, source, source_id, now_iso()),
        )
    return job_id


def finish_job(
    job_id: str,
    *,
    status: str,
    documents: int = 0,
    chunks: int = 0,
    error: str | None = None,
) -> None:
    with db() as conn:
        conn.execute(
            "update index_jobs set status = ?, documents = ?, chunks = ?, error = ?, finished_at = ? where id = ?",
            (status, documents, chunks, error, now_iso(), job_id),
        )


def update_attachment_status(
    attachment: dict[str, Any],
    *,
    status: str,
    chunks: int = 0,
    error: str | None = None,
) -> None:
    with db() as conn:
        conn.execute(
            """
            insert into attachment_index (
              id, document_id, filename, content_type, size, status, chunks,
              error, updated_at, indexed_at
            )
            values (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            on conflict(id) do update set
              document_id = excluded.document_id,
              filename = excluded.filename,
              content_type = excluded.content_type,
              size = excluded.size,
              status = excluded.status,
              chunks = excluded.chunks,
              error = excluded.error,
              updated_at = excluded.updated_at,
              indexed_at = excluded.indexed_at
            """,
            (
                attachment["id"],
                attachment["documentId"],
                Path(attachment["key"]).name,
                attachment["contentType"],
                attachment["size"],
                status,
                chunks,
                error,
                attachment["updatedAt"],
                now_iso() if status in {"indexed", "needs_ocr", "error"} else None,
            ),
        )


def update_document_status(
    document: dict[str, Any],
    *,
    status: str,
    chunks: int = 0,
    error: str | None = None,
) -> None:
    with db() as conn:
        conn.execute(
            """
            insert into document_index (
              id, title, collection, status, chunks, error, updated_at, indexed_at
            )
            values (?, ?, ?, ?, ?, ?, ?, ?)
            on conflict(id) do update set
              title = excluded.title,
              collection = excluded.collection,
              status = excluded.status,
              chunks = excluded.chunks,
              error = excluded.error,
              updated_at = excluded.updated_at,
              indexed_at = excluded.indexed_at
            """,
            (
                document["id"],
                repair_mojibake(document.get("title")),
                repair_mojibake(document.get("collection")),
                status,
                chunks,
                error,
                document["updatedAt"],
                now_iso() if status in {"indexed", "skipped", "error"} else None,
            ),
        )


def document_needs_index(document: dict[str, Any]) -> bool:
    with db() as conn:
        row = conn.execute("select status, updated_at from document_index where id = ?", (document["id"],)).fetchone()
    return not row or row["updated_at"] != document["updatedAt"] or row["status"] in {"running"}


def attachment_needs_index(attachment: dict[str, Any]) -> bool:
    with db() as conn:
        row = conn.execute("select status, updated_at from attachment_index where id = ?", (attachment["id"],)).fetchone()
    return not row or row["updated_at"] != attachment["updatedAt"] or row["status"] in {"indexing"}


async def index_outline_document(document: dict[str, Any]) -> int:
    title = clean_text(document["title"])
    full_text = f"{title}\n\n{clean_text(document['text'])}".strip()
    if useful_text_length(full_text) < MIN_OUTLINE_TEXT_CHARS:
        delete_source_points("outline", document["id"])
        update_document_status(document, status="skipped", chunks=0, error="Not enough text to index.")
        return 0
    update_document_status(document, status="running")
    chunks = await index_chunks(
        source_type="outline",
        source_id=document["id"],
        title=title,
        url=document_url(title, document["urlId"]),
        text=full_text,
        extra_payload={
            "documentId": document["id"],
            "collection": document["collection"],
            "collectionId": document["collectionId"],
            "updatedAt": document["updatedAt"],
        },
    )
    update_document_status(document, status="indexed", chunks=chunks)
    return chunks


async def index_outline_attachment(attachment: dict[str, Any]) -> int:
    path = OUTLINE_FILES_DIR / attachment["key"]
    if not path.exists():
        update_attachment_status(attachment, status="error", error="Stored file is missing.")
        return 0
    update_attachment_status(attachment, status="indexing")
    try:
        text = extract_file_text(path)
    except Exception as exc:
        update_attachment_status(attachment, status="error", error=str(exc))
        return 0
    if not text:
        delete_source_points("outline_attachment", attachment["id"])
        update_attachment_status(attachment, status="needs_ocr", error="No text layer extracted. OCR is required.")
        return 0
    filename = Path(attachment["key"]).name
    title = f"{attachment['documentTitle']} � {filename}"
    document_page_url = document_url(attachment["documentTitle"], attachment["urlId"])
    chunks = await index_chunks(
        source_type="outline_attachment",
        source_id=attachment["id"],
        title=title,
        url=document_page_url,
        text=f"{title}\n\n{text}",
        extra_payload={
            "documentId": attachment["documentId"],
            "documentUrl": document_page_url,
            "attachmentId": attachment["id"],
            "collection": attachment["collection"],
            "collectionId": attachment.get("collectionId"),
            "filename": filename,
            "contentType": attachment["contentType"],
            "updatedAt": attachment["updatedAt"],
        },
    )
    update_attachment_status(attachment, status="indexed", chunks=chunks)
    return chunks


async def run_incremental_index() -> dict[str, int]:
    if index_lock.locked():
        return {"documents": 0, "chunks": 0, "skipped": 1}
    async with index_lock:
        pending_documents = [document for document in fetch_documents() if document_needs_index(document)]
        pending_attachments = [attachment for attachment in fetch_indexable_attachments() if attachment_needs_index(attachment)]
        if not pending_documents and not pending_attachments:
            return {"documents": 0, "chunks": 0, "skipped": 0}
        job_id = create_job("outline-auto")
        indexed_documents = 0
        indexed_chunks = 0
        try:
            for document in pending_documents:
                try:
                    indexed_chunks += await index_outline_document(document)
                    indexed_documents += 1
                except Exception as exc:
                    update_document_status(document, status="error", error=str(exc))
            for attachment in pending_attachments:
                indexed_chunks += await index_outline_attachment(attachment)
                indexed_documents += 1
            finish_job(job_id, status="done", documents=indexed_documents, chunks=indexed_chunks)
            return {"documents": indexed_documents, "chunks": indexed_chunks, "skipped": 0}
        except Exception as exc:
            finish_job(job_id, status="error", documents=indexed_documents, chunks=indexed_chunks, error=str(exc))
            raise


async def run_outline_index() -> dict[str, int]:
    if index_lock.locked():
        return {"documents": 0, "chunks": 0, "skipped": 1}

    async with index_lock:
        job_id = create_job("outline")
        try:
            documents = fetch_documents()
            attachments = fetch_indexable_attachments()
            indexed_chunks = 0
            delete_source_type_points("outline")
            delete_source_type_points("outline_pdf")
            delete_source_type_points("outline_attachment")
            for document in documents:
                title = clean_text(document["title"])
                full_text = f"{title}\n\n{clean_text(document['text'])}".strip()
                if useful_text_length(full_text) < MIN_OUTLINE_TEXT_CHARS:
                    continue
                chunks = await index_chunks(
                    source_type="outline",
                    source_id=document["id"],
                    title=title,
                    url=document_url(title, document["urlId"]),
                    text=full_text,
                    extra_payload={
                        "documentId": document["id"],
                        "collection": document["collection"],
                        "collectionId": document["collectionId"],
                        "updatedAt": document["updatedAt"],
                    },
                )
                indexed_chunks += chunks
            for attachment in attachments:
                path = OUTLINE_FILES_DIR / attachment["key"]
                if not path.exists():
                    update_attachment_status(
                        attachment,
                        status="error",
                        error="Stored file is missing.",
                    )
                    continue
                update_attachment_status(attachment, status="indexing")
                try:
                    text = extract_file_text(path)
                except Exception as exc:
                    update_attachment_status(
                        attachment,
                        status="error",
                        error=str(exc),
                    )
                    continue
                if not text:
                    update_attachment_status(
                        attachment,
                        status="needs_ocr",
                        error="No text layer extracted. OCR is required.",
                    )
                    continue
                filename = Path(attachment["key"]).name
                title = f"{attachment['documentTitle']} · {filename}"
                document_page_url = document_url(attachment["documentTitle"], attachment["urlId"])
                chunks = await index_chunks(
                    source_type="outline_attachment",
                    source_id=attachment["id"],
                    title=title,
                    url=document_page_url,
                    text=f"{title}\n\n{text}",
                    extra_payload={
                        "documentId": attachment["documentId"],
                        "documentUrl": document_page_url,
                        "attachmentId": attachment["id"],
                        "collection": attachment["collection"],
                        "collectionId": attachment.get("collectionId"),
                        "filename": filename,
                        "contentType": attachment["contentType"],
                        "updatedAt": attachment["updatedAt"],
                    },
                )
                update_attachment_status(attachment, status="indexed", chunks=chunks)
                indexed_chunks += chunks
            indexed_documents = len(documents) + len(attachments)
            finish_job(job_id, status="done", documents=indexed_documents, chunks=indexed_chunks)
            return {"documents": indexed_documents, "chunks": indexed_chunks, "skipped": 0}
        except Exception as exc:
            finish_job(job_id, status="error", error=str(exc))
            raise


async def auto_index_loop() -> None:
    await asyncio.sleep(5)
    while True:
        with suppress(Exception):
            await run_incremental_index()
        await asyncio.sleep(FAST_INDEX_INTERVAL_SECONDS)


@app.get("/health")
async def health() -> dict[str, str]:
    return {"status": "ok"}


async def probe_ollama() -> dict[str, Any]:
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            response = await client.get(f"{OLLAMA_URL}/api/tags")
        response.raise_for_status()
        models = response.json().get("models", [])
        return {"ok": True, "models": [item.get("name") or item.get("model") for item in models]}
    except Exception as exc:
        return {"ok": False, "error": str(exc)}


def qdrant_status() -> dict[str, Any]:
    try:
        collections = [item.name for item in qdrant.get_collections().collections]
        if COLLECTION not in collections:
            return {"ok": True, "collection": COLLECTION, "exists": False, "points": 0}
        info = qdrant.get_collection(COLLECTION)
        return {
            "ok": True,
            "collection": COLLECTION,
            "exists": True,
            "points": int(info.points_count or 0),
            "indexedVectors": int(info.indexed_vectors_count or 0),
            "status": str(info.status),
        }
    except Exception as exc:
        return {"ok": False, "collection": COLLECTION, "error": str(exc)}


def normalize_status(value: str | None) -> str:
    if value in {"done", "running", "error", "indexed", "needs_ocr"}:
        return value
    if not value:
        return "unknown"
    lowered = value.lower()
    if "error" in lowered or "????" in lowered:
        return "error"
    if "index" in lowered or "??????" in lowered:
        return "running"
    if "done" in lowered or "?????" in lowered:
        return "done"
    return value


def normalize_job_status(item: dict[str, Any]) -> str:
    status = normalize_status(item.get("status"))
    if status in {"done", "running", "error"}:
        return status
    if item.get("error"):
        return "error"
    if item.get("finished_at"):
        return "done"
    return "running"

def latest_job() -> dict[str, Any] | None:
    with db() as conn:
        row = conn.execute(
            """
            select id, source, source_id, status, documents, chunks, error, created_at, finished_at
            from index_jobs order by created_at desc limit 1
            """
        ).fetchone()
    if not row:
        return None
    item = dict(row)
    item["status"] = normalize_job_status(item)
    return item


@app.get("/status")
async def status() -> dict[str, Any]:
    collections = fetch_collections()
    return {
        "status": "ok",
        "chatModel": CHAT_MODEL,
        "embedModel": EMBED_MODEL,
        "autoIndexIntervalSeconds": AUTO_INDEX_INTERVAL_SECONDS,
        "ollama": await probe_ollama(),
        "qdrant": qdrant_status(),
        "collections": len(collections),
        "latestJob": latest_job(),
    }


@app.get("/", response_class=HTMLResponse)
async def index_page() -> str:
    return HTML_PAGE


@app.post("/index")
async def index_outline() -> dict[str, int]:
    return await run_outline_index()


@app.post("/files/upload")
async def upload_file(file: UploadFile = File(...), collection_id: str | None = Form(None)) -> dict[str, Any]:
    if not file.filename:
        raise HTTPException(status_code=400, detail="Р вЂ”Р В°Р С–РЎР‚РЎС“Р В·Р С‘РЎвЂљР Вµ РЎвЂћР В°Р в„–Р В».")
    extension = Path(file.filename).suffix.lower()
    if extension in LEGACY_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail="Р РЋРЎвЂљР В°РЎР‚РЎвЂ№Р Вµ РЎвЂћР С•РЎР‚Р СР В°РЎвЂљРЎвЂ№ .doc Р С‘ .xls Р С—Р С•Р С”Р В° Р Р…Р Вµ Р С‘Р Р…Р Т‘Р ВµР С”РЎРѓР С‘РЎР‚РЎС“РЎР‹РЎвЂљРЎРѓРЎРЏ. Р РЋР С•РЎвЂ¦РЎР‚Р В°Р Р…Р С‘РЎвЂљР Вµ РЎвЂћР В°Р в„–Р В» Р С”Р В°Р С” .docx Р С‘Р В»Р С‘ .xlsx Р С‘ Р В·Р В°Р С–РЎР‚РЎС“Р В·Р С‘РЎвЂљР Вµ Р В·Р В°Р Р…Р С•Р Р†Р С•.",
        )
    if extension not in SUPPORTED_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Р СџР С•Р Т‘Р Т‘Р ВµРЎР‚Р В¶Р С‘Р Р†Р В°РЎР‹РЎвЂљРЎРѓРЎРЏ РЎвЂљР С•Р В»РЎРЉР С”Р С• PDF, DOCX Р С‘ XLSX.")
    file_id = str(uuid.uuid4())
    safe_name = re.sub(r"[^A-Za-z0-9Р С’-Р Р‡Р В°-РЎРЏР РѓРЎвЂ._ -]+", "_", file.filename).strip() or f"document{extension}"
    path = FILES_DIR / f"{file_id}{extension}"
    collection = get_collection(collection_id)
    size = 0
    with path.open("wb") as output:
        while chunk := await file.read(1024 * 1024):
            size += len(chunk)
            output.write(chunk)
    with db() as conn:
        conn.execute(
            """
            insert into files (id, filename, title, content_type, size, path, status, collection_id, collection_name, created_at)
            values (?, ?, ?, ?, ?, ?, 'Р В·Р В°Р С–РЎР‚РЎС“Р В¶Р ВµР Р…', ?, ?, ?)
            """,
            (
                file_id,
                safe_name,
                Path(safe_name).stem,
                file.content_type or MEDIA_TYPES.get(extension),
                size,
                str(path),
                collection["id"] if collection else None,
                collection["name"] if collection else None,
                now_iso(),
            ),
        )
    return await index_file(file_id)


@app.post("/files/convert")
async def convert_file(file: UploadFile = File(...)) -> ConvertResponse:
    if not file.filename:
        raise HTTPException(status_code=400, detail="Р вЂ”Р В°Р С–РЎР‚РЎС“Р В·Р С‘РЎвЂљР Вµ РЎвЂћР В°Р в„–Р В».")
    extension = Path(file.filename).suffix.lower()
    if extension in LEGACY_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail="Р РЋРЎвЂљР В°РЎР‚РЎвЂ№Р Вµ РЎвЂћР С•РЎР‚Р СР В°РЎвЂљРЎвЂ№ .doc Р С‘ .xls Р С—Р С•Р С”Р В° Р Р…Р Вµ Р С‘Р СР С—Р С•РЎР‚РЎвЂљР С‘РЎР‚РЎС“РЎР‹РЎвЂљРЎРѓРЎРЏ. Р РЋР С•РЎвЂ¦РЎР‚Р В°Р Р…Р С‘РЎвЂљР Вµ РЎвЂћР В°Р в„–Р В» Р С”Р В°Р С” .docx Р С‘Р В»Р С‘ .xlsx Р С‘ Р В·Р В°Р С–РЎР‚РЎС“Р В·Р С‘РЎвЂљР Вµ Р В·Р В°Р Р…Р С•Р Р†Р С•.",
        )
    if extension not in SUPPORTED_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Р СџР С•Р Т‘Р Т‘Р ВµРЎР‚Р В¶Р С‘Р Р†Р В°РЎР‹РЎвЂљРЎРѓРЎРЏ РЎвЂљР С•Р В»РЎРЉР С”Р С• PDF, DOCX Р С‘ XLSX.")

    safe_name = re.sub(r"[^A-Za-z0-9Р С’-Р Р‡Р В°-РЎРЏР РѓРЎвЂ._ -]+", "_", file.filename).strip() or f"document{extension}"
    temp_path = DATA_DIR / f"convert-{uuid.uuid4()}{extension}"
    try:
        with temp_path.open("wb") as output:
            while chunk := await file.read(1024 * 1024):
                output.write(chunk)
        text = extract_file_text(temp_path)
        if not text:
            raise ValueError("Р СњР Вµ РЎС“Р Т‘Р В°Р В»Р С•РЎРѓРЎРЉ Р С‘Р В·Р Р†Р В»Р ВµРЎвЂЎРЎРЉ РЎвЂљР ВµР С”РЎРѓРЎвЂљ. Р вЂќР В»РЎРЏ PDF-РЎРѓР С”Р В°Р Р…Р С•Р Р† Р Р…РЎС“Р В¶Р ВµР Р… OCR.")
        return ConvertResponse(title=Path(safe_name).stem, text=text)
    except ValueError as exc:
        raise HTTPException(status_code=422, detail=str(exc))
    finally:
        if temp_path.exists():
            temp_path.unlink()


@app.post("/files/{file_id}/index")
async def index_file(file_id: str) -> dict[str, Any]:
    job_id = create_job("file", file_id)
    with db() as conn:
        row = conn.execute("select * from files where id = ?", (file_id,)).fetchone()
    if not row:
        finish_job(job_id, status="error", error="File not found")
        raise HTTPException(status_code=404, detail="Р В¤Р В°Р в„–Р В» Р Р…Р Вµ Р Р…Р В°Р в„–Р Т‘Р ВµР Р….")
    try:
        with db() as conn:
            conn.execute("update files set status = 'running', error = null where id = ?", (file_id,))
        path = Path(row["path"])
        text = extract_file_text(path)
        if not text:
            raise ValueError("Р СњР Вµ РЎС“Р Т‘Р В°Р В»Р С•РЎРѓРЎРЉ Р С‘Р В·Р Р†Р В»Р ВµРЎвЂЎРЎРЉ РЎвЂљР ВµР С”РЎРѓРЎвЂљ. Р вЂќР В»РЎРЏ PDF-РЎРѓР С”Р В°Р Р…Р С•Р Р† Р Р…РЎС“Р В¶Р ВµР Р… OCR, Р В° Word/Excel Р Т‘Р С•Р В»Р В¶Р Р…РЎвЂ№ Р В±РЎвЂ№РЎвЂљРЎРЉ Р Р† РЎвЂћР С•РЎР‚Р СР В°РЎвЂљР Вµ .docx/.xlsx.")
        chunks = await index_chunks(
            source_type="file",
            source_id=file_id,
            title=row["title"],
            url=f"{OUTLINE_URL}/ai/files/{file_id}/download",
            text=f"{row['title']}\n\n{text}",
            extra_payload={
                "filename": row["filename"],
                "contentType": row["content_type"],
                "fileType": path.suffix.lower().lstrip("."),
                "collection": row["collection_name"],
                "collectionId": row["collection_id"],
            },
        )
        with db() as conn:
            conn.execute(
                "update files set status = 'done', chunks = ?, indexed_at = ?, error = null where id = ?",
                (chunks, now_iso(), file_id),
            )
        finish_job(job_id, status="done", documents=1, chunks=chunks)
        return {"id": file_id, "filename": row["filename"], "status": "done", "chunks": chunks}
    except Exception as exc:
        message = str(exc)
        with db() as conn:
            conn.execute("update files set status = 'error', error = ? where id = ?", (message, file_id))
        finish_job(job_id, status="error", error=message)
        raise HTTPException(status_code=422, detail=message)


@app.get("/files")
async def list_files() -> dict[str, Any]:
    with db() as conn:
        rows = conn.execute(
            """
            select id, filename, title, size, status, chunks, error, collection_id, collection_name, created_at, indexed_at
            from files order by created_at desc
            """
        ).fetchall()
    return {"files": [dict(row) for row in rows]}


@app.get("/collections")
async def list_collections() -> dict[str, Any]:
    return {"collections": fetch_collections()}


@app.get("/files/{file_id}/download")
async def download_file(file_id: str) -> FileResponse:
    with db() as conn:
        row = conn.execute("select * from files where id = ?", (file_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Р В¤Р В°Р в„–Р В» Р Р…Р Вµ Р Р…Р В°Р в„–Р Т‘Р ВµР Р….")
    media_type = row["content_type"] or MEDIA_TYPES.get(Path(row["path"]).suffix.lower()) or "application/octet-stream"
    return FileResponse(path=row["path"], media_type=media_type, filename=row["filename"])


@app.delete("/files/{file_id}")
async def delete_file(file_id: str) -> dict[str, str]:
    with db() as conn:
        row = conn.execute("select * from files where id = ?", (file_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Р В¤Р В°Р в„–Р В» Р Р…Р Вµ Р Р…Р В°Р в„–Р Т‘Р ВµР Р….")

    delete_source_points("file", file_id)
    delete_source_points("pdf", file_id)
    path = Path(row["path"])
    if path.exists():
        path.unlink()
    with db() as conn:
        conn.execute("delete from files where id = ?", (file_id,))
    return {"status": "РЎС“Р Т‘Р В°Р В»Р ВµР Р…Р С•"}


@app.get("/document-status")
async def document_status() -> dict[str, Any]:
    documents = fetch_documents()
    attachments = fetch_indexable_attachments()
    with db() as conn:
        doc_rows = {row["id"]: dict(row) for row in conn.execute("select * from document_index").fetchall()}
        att_rows = {row["id"]: dict(row) for row in conn.execute("select * from attachment_index").fetchall()}
    items: list[dict[str, Any]] = []
    for document in documents:
        row = doc_rows.get(document["id"], {})
        status = row.get("status") or "pending"
        if row.get("updated_at") and row.get("updated_at") != document["updatedAt"]:
            status = "pending"
        items.append({
            "id": document["id"],
            "type": "page",
            "title": repair_mojibake(document.get("title")),
            "collection": repair_mojibake(document.get("collection")),
            "status": status,
            "chunks": row.get("chunks", 0),
            "error": row.get("error"),
            "updatedAt": document["updatedAt"],
            "indexedAt": row.get("indexed_at"),
        })
    for attachment in attachments:
        row = att_rows.get(attachment["id"], {})
        status = row.get("status") or "pending"
        if row.get("updated_at") and row.get("updated_at") != attachment["updatedAt"]:
            status = "pending"
        items.append({
            "id": attachment["id"],
            "documentId": attachment["documentId"],
            "type": "file",
            "title": repair_mojibake(Path(attachment["key"]).name),
            "documentTitle": repair_mojibake(attachment.get("documentTitle")),
            "collection": repair_mojibake(attachment.get("collection")),
            "status": status,
            "chunks": row.get("chunks", 0),
            "error": row.get("error"),
            "updatedAt": attachment["updatedAt"],
            "indexedAt": row.get("indexed_at"),
        })
    order = {"indexing": 0, "running": 0, "pending": 1, "error": 2, "needs_ocr": 3, "indexed": 4, "skipped": 5}
    items.sort(key=lambda item: (order.get(item["status"], 9), item.get("title") or ""))
    summary = {"total": len(items), "indexed": 0, "pending": 0, "indexing": 0, "error": 0, "needs_ocr": 0, "skipped": 0}
    for item in items:
        status = item["status"]
        if status in {"indexed", "done"}:
            summary["indexed"] += 1
        elif status in {"indexing", "running"}:
            summary["indexing"] += 1
        elif status in {"error"}:
            summary["error"] += 1
        elif status in {"needs_ocr"}:
            summary["needs_ocr"] += 1
        elif status in {"skipped"}:
            summary["skipped"] += 1
        else:
            summary["pending"] += 1
    return {"summary": summary, "items": items[:80]}


@app.get("/jobs")
async def list_jobs() -> dict[str, Any]:
    with db() as conn:
        rows = conn.execute(
            """
            select id, source, source_id, status, documents, chunks, error, created_at, finished_at
            from index_jobs order by created_at desc limit 50
            """
        ).fetchall()
    jobs = []
    for row in rows:
        item = dict(row)
        item["status"] = normalize_job_status(item)
        jobs.append(item)
    return {"jobs": jobs}


@app.post("/search")
async def search(request: SearchRequest) -> dict[str, Any]:
    if len(clean_text(request.query)) < 3:
        return {"results": []}
    vector = await embed(expand_query(request.query))
    await ensure_collection(len(vector))
    fetch_limit = max(request.limit * 12, request.limit, 50)
    raw_results = qdrant.search(
        collection_name=COLLECTION,
        query_vector=vector,
        limit=fetch_limit,
        with_payload=True,
    )
    allowed_document_ids = set(request.allowedDocumentIds or [])
    results: list[dict[str, Any]] = []
    for result in raw_results:
        payload = result.payload or {}
        document_id = payload.get("documentId") or payload.get("sourceId")
        if request.collectionId and payload.get("collectionId") != request.collectionId:
            continue
        if request.allowedDocumentIds is not None and document_id not in allowed_document_ids:
            continue
        title_boost = title_match_score(request.query, payload)
        score = float(result.score) + title_boost
        if score < MIN_SEARCH_SCORE and not title_boost:
            continue
        source_type = payload.get("sourceType")
        source_url = payload.get("url")
        document_page_url = payload.get("documentUrl")
        if source_type in {"outline", "outline_attachment"}:
            document_page_url = lookup_document_url(document_id) or document_page_url
            source_url = document_page_url or source_url
        results.append(
            {
                "score": score,
                "rawScore": float(result.score),
                "titleMatch": bool(title_boost),
                "sourceType": source_type,
                "documentId": document_id,
                "attachmentId": payload.get("attachmentId"),
                "documentUrl": document_page_url,
                "title": repair_mojibake(payload.get("title")),
                "collection": repair_mojibake(payload.get("collection")),
                "collectionId": payload.get("collectionId"),
                "filename": repair_mojibake(payload.get("filename")),
                "url": source_url,
                "chunk": payload.get("chunk"),
            }
        )
    results.sort(key=lambda item: item.get("score", 0), reverse=True)
    return {"results": results[: request.limit]}


def conversation_context(history: list[ChatHistoryMessage], limit: int = 8) -> str:
    lines: list[str] = []
    for item in history[-limit:]:
        role = item.role.lower().strip()
        if role not in {"user", "assistant"}:
            continue
        content = clean_text(item.content)[:700]
        if not content:
            continue
        label = "User" if role == "user" else "Assistant"
        lines.append(f"{label}: {content}")
    return "\n".join(lines)


@app.post("/chat")
async def chat(request: ChatRequest) -> dict[str, Any]:
    no_info = "\u0412 \u0431\u0430\u0437\u0435 \u0437\u043d\u0430\u043d\u0438\u0439 \u043d\u0435 \u043d\u0430\u0439\u0434\u0435\u043d\u043e \u0434\u043e\u0441\u0442\u0430\u0442\u043e\u0447\u043d\u043e \u0438\u043d\u0444\u043e\u0440\u043c\u0430\u0446\u0438\u0438 \u0434\u043b\u044f \u043e\u0442\u0432\u0435\u0442\u0430."
    if len(clean_text(request.message)) < 3:
        return {"answer": "\u0423\u0442\u043e\u0447\u043d\u0438\u0442\u0435 \u0432\u043e\u043f\u0440\u043e\u0441: \u0437\u0430\u043f\u0440\u043e\u0441 \u0441\u043b\u0438\u0448\u043a\u043e\u043c \u043a\u043e\u0440\u043e\u0442\u043a\u0438\u0439 \u0434\u043b\u044f \u043f\u043e\u0438\u0441\u043a\u0430 \u043f\u043e \u0431\u0430\u0437\u0435 \u0437\u043d\u0430\u043d\u0438\u0439.", "sources": []}

    search_results = await search(
        SearchRequest(
            query=request.message,
            limit=max(request.limit * 6, 18),
            collectionId=request.collectionId,
            allowedDocumentIds=request.allowedDocumentIds,
        )
    )
    results = search_results["results"]
    if not results:
        return {"answer": no_info, "sources": []}

    if is_inventory_request(request.message):
        inventory_limit = min(max(request.limit, 8), 12)
        inventory_results = context_results(results, limit=inventory_limit) or results[:inventory_limit]
        return {
            "answer": build_sources_fallback(inventory_results, "\u0412 \u0431\u0430\u0437\u0435 \u0437\u043d\u0430\u043d\u0438\u0439 \u043d\u0430\u0439\u0434\u0435\u043d\u044b \u043c\u0430\u0442\u0435\u0440\u0438\u0430\u043b\u044b:"),
            "sources": readable_sources(inventory_results),
        }

    selected_results = focused_results(results, limit=min(max(request.limit, 1), 4))
    if not selected_results:
        return {"answer": no_info, "sources": []}

    context_items = [
        f"Source: {item.get('sourceType')}\nTitle: {item.get('title')}\nText: {item.get('chunk')}"
        for item in selected_results
    ]
    context = "\n\n".join(context_items)
    history_context = conversation_context(request.history)
    history_block = f"Conversation history:\n{history_context}\n\n" if history_context else ""
    prompt = (
        "Answer in Russian using only the Outline and uploaded-file context below. "
        "Use the conversation history only to understand follow-up questions, not as a factual source. "
        "If the context is insufficient, say exactly: " + no_info + " "
        "Do not add facts that are absent from the context. Do not include URLs in the answer; sources are shown separately. "
        "Keep the answer concise and practical.\n\n"
        f"{history_block}"
        f"Context:\n{context}\n\n"
        f"Question: {request.message}"
    )
    try:
        async with httpx.AsyncClient(timeout=420) as client:
            response = await client.post(
                f"{OLLAMA_URL}/api/chat",
                json={
                    "model": CHAT_MODEL,
                    "messages": [{"role": "user", "content": prompt}],
                    "stream": False,
                    "think": False,
                    "keep_alive": CHAT_KEEP_ALIVE,
                    "options": {"num_ctx": 2048, "num_predict": 256, "temperature": 0.2},
                },
            )
        if response.status_code >= 400:
            raise httpx.HTTPStatusError(f"Ollama returned {response.status_code}", request=response.request, response=response)
        answer = readable_answer(response.json().get("message", {}).get("content", ""))
        if not answer:
            answer = build_sources_fallback(selected_results, "\u041c\u043e\u0434\u0435\u043b\u044c \u0432\u0440\u0435\u043c\u0435\u043d\u043d\u043e \u043d\u0435 \u0432\u0435\u0440\u043d\u0443\u043b\u0430 \u0442\u0435\u043a\u0441\u0442, \u043d\u043e \u0438\u043d\u0434\u0435\u043a\u0441 \u043d\u0430\u0448\u0435\u043b \u043c\u0430\u0442\u0435\u0440\u0438\u0430\u043b\u044b:")
        if clean_text(answer) == no_info:
            return {"answer": answer, "sources": []}
        return {"answer": answer, "sources": readable_sources(selected_results)}
    except Exception:
        return {"answer": build_sources_fallback(selected_results, "\u041d\u0435 \u0443\u0434\u0430\u043b\u043e\u0441\u044c \u043f\u043e\u043b\u0443\u0447\u0438\u0442\u044c \u043e\u0442\u0432\u0435\u0442 \u043e\u0442 \u043c\u043e\u0434\u0435\u043b\u0438, \u043d\u043e \u0438\u043d\u0434\u0435\u043a\u0441 \u043d\u0430\u0448\u0435\u043b \u043c\u0430\u0442\u0435\u0440\u0438\u0430\u043b\u044b:"), "sources": readable_sources(selected_results)}


def is_inventory_request(message: str) -> bool:
    normalized = clean_text(message).lower()
    phrases = (
        "что есть в базе",
        "что есть в базе знаний",
        "что в базе",
        "что загружено",
        "какие документы есть",
        "какие материалы есть",
        "перечисли документы",
        "список документов",
        "what is in the knowledge base",
        "list documents",
    )
    return any(phrase in normalized for phrase in phrases)


def build_sources_fallback(results: list[dict[str, Any]], prefix: str) -> str:
    unique_titles: list[str] = []
    seen: set[str] = set()
    for item in results:
        title = clean_text(item.get("title") or item.get("filename") or "Document")
        collection = clean_text(item.get("collection"))
        source_type = item.get("sourceType")
        suffix = "file" if source_type in {"file", "pdf", "outline_pdf", "outline_attachment"} else "page"
        label = f"{title} - {collection} ({suffix})" if collection else f"{title} ({suffix})"
        if label in seen:
            continue
        seen.add(label)
        unique_titles.append(label)
    if not unique_titles:
        return "\u0412 \u0431\u0430\u0437\u0435 \u0437\u043d\u0430\u043d\u0438\u0439 \u043d\u0435 \u043d\u0430\u0439\u0434\u0435\u043d\u043e \u0434\u043e\u0441\u0442\u0430\u0442\u043e\u0447\u043d\u043e \u0438\u043d\u0444\u043e\u0440\u043c\u0430\u0446\u0438\u0438 \u0434\u043b\u044f \u043e\u0442\u0432\u0435\u0442\u0430."
    lines = "\n".join(f"- {title}" for title in unique_titles[:12])
    return f"{prefix}\n{lines}"


HTML_PAGE = r"""
<!doctype html>
<html lang="ru">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>AI &#1087;&#1086;&#1080;&#1089;&#1082; - Outline</title>
  <style>
    :root{--bg:#fff;--text:#111827;--muted:#687386;--line:#dfe5ee;--soft:#edf1f6;--panel:#fbfcfe;--primary:#0f6fec;--hover:#0b5ed7;--danger:#9d3128;color-scheme:light}
    *{box-sizing:border-box}html,body{margin:0;min-height:100%;background:var(--bg);color:var(--text);font-family:ui-sans-serif,system-ui,-apple-system,BlinkMacSystemFont,"Segoe UI",sans-serif;font-size:15px;line-height:1.5}button,textarea{font:inherit}
    .page{width:min(1120px,calc(100vw - 56px));min-height:100vh;margin:0 auto;padding:34px 0 28px;display:grid;grid-template-rows:auto minmax(300px,1fr) auto;gap:16px}
    header{display:flex;justify-content:space-between;gap:18px;align-items:flex-start;border-bottom:1px solid var(--line);padding-bottom:16px}h1{margin:0;font-size:30px;line-height:1.18;font-weight:700;letter-spacing:0}.subtitle{margin-top:7px;color:var(--muted);max-width:760px}.toolbar{display:flex;gap:10px;align-items:center;flex-wrap:wrap}.status{min-height:34px;padding:7px 10px;border:1px solid var(--line);border-radius:6px;color:var(--muted);background:#f8fafc;font-size:13px;white-space:nowrap}
    button,.nav-button{min-height:40px;border:1px solid var(--primary);border-radius:6px;background:var(--primary);color:#fff;padding:0 16px;font-weight:650;cursor:pointer;box-shadow:0 1px 2px rgba(15,111,236,.18);display:inline-flex;align-items:center;text-decoration:none}button:hover,.nav-button:hover{background:var(--hover)}button:disabled{opacity:.58;cursor:wait}.secondary{border-color:#c7d0dd;background:#fff;color:#2f3b4a;box-shadow:none}.secondary:hover{background:#f8fafc}
    .doc-status{border:1px solid var(--line);border-radius:8px;background:#fff;padding:12px}.doc-status-head{display:flex;justify-content:space-between;gap:12px;align-items:center;margin-bottom:8px}.doc-status-title{font-weight:700}.doc-status-summary{font-size:13px;color:var(--muted)}.doc-status-list{display:grid;gap:6px;max-height:190px;overflow:auto}.doc-row{display:grid;grid-template-columns:minmax(0,1fr) auto;gap:10px;align-items:center;border-top:1px solid var(--soft);padding:7px 0}.doc-row:first-child{border-top:0}.doc-name{overflow:hidden;text-overflow:ellipsis;white-space:nowrap}.badge{border:1px solid var(--line);border-radius:999px;padding:2px 8px;font-size:12px;font-weight:700;white-space:nowrap}.badge.error{color:var(--danger);border-color:var(--danger)}.badge.pending{color:#8a5a00;border-color:#d7a62f}.badge.indexed{color:#17633a;border-color:#5ab581}.badge.needs_ocr{color:#7a3b00;border-color:#d89443}
    .chat-panel{min-width:0;border:1px solid var(--line);border-radius:8px;background:#fff;display:grid;grid-template-rows:auto minmax(280px,1fr);overflow:hidden}.panel-head{padding:16px 20px;border-bottom:1px solid var(--line);background:var(--panel)}h2{margin:0 0 5px;font-size:18px;line-height:1.3}.small{color:var(--muted);font-size:13px}.panel-title-row{display:flex;align-items:flex-start;justify-content:space-between;gap:12px}.panel-title-row .secondary{padding:8px 10px;font-size:13px;white-space:nowrap}
    #messages{padding:22px;overflow:auto;min-height:360px;max-height:calc(100vh - 330px)}.message{max-width:860px;margin:0 0 18px;padding:0 0 18px;border-bottom:1px solid var(--soft);white-space:pre-wrap;overflow-wrap:anywhere}.role{display:block;margin-bottom:6px;color:#344255;font-size:12px;font-weight:700;text-transform:uppercase;letter-spacing:0}.sources{display:grid;gap:8px;margin-top:12px}.sources a{width:fit-content;max-width:100%;color:var(--primary);text-decoration:none;font-size:13px;font-weight:600;overflow-wrap:anywhere}.sources a:hover{text-decoration:underline}
    form.chat{display:grid;grid-template-columns:minmax(0,1fr) auto;gap:12px;align-items:end;border:1px solid var(--line);border-radius:8px;padding:12px;background:#fff}textarea{width:100%;min-height:52px;max-height:180px;resize:vertical;border:1px solid #c7d0dd;border-radius:6px;padding:12px;color:var(--text);background:#fff;outline:none}textarea:focus,button:focus-visible{outline:2px solid rgba(15,111,236,.26);outline-offset:2px}
    @media(max-width:820px){.page{width:calc(100vw - 28px);padding-top:22px}header,form.chat{display:grid;grid-template-columns:1fr}.toolbar{justify-content:flex-start}.status{white-space:normal}.panel-title-row{align-items:flex-start}#messages{max-height:none}}
  </style>
</head>
<body>
  <main class="page">
    <header>
      <div>
        <h1>AI &#1087;&#1086;&#1080;&#1089;&#1082; Outline</h1>
        <div class="subtitle">&#1063;&#1072;&#1090; &#1086;&#1090;&#1074;&#1077;&#1095;&#1072;&#1077;&#1090; &#1087;&#1086; &#1076;&#1086;&#1082;&#1091;&#1084;&#1077;&#1085;&#1090;&#1072;&#1084; &#1080; &#1074;&#1083;&#1086;&#1078;&#1077;&#1085;&#1080;&#1103;&#1084; Outline. &#1048;&#1089;&#1090;&#1086;&#1095;&#1085;&#1080;&#1082;&#1080; &#1087;&#1086;&#1082;&#1072;&#1079;&#1099;&#1074;&#1072;&#1102;&#1090;&#1089;&#1103; &#1087;&#1086;&#1076; &#1086;&#1090;&#1074;&#1077;&#1090;&#1086;&#1084;.</div>
      </div>
      <div class="toolbar"><a class="nav-button secondary" href="/outline/" target="_top">Outline</a></div>
    </header>
    <section class="doc-status"><div class="doc-status-head"><div class="doc-status-title">&#1057;&#1090;&#1072;&#1090;&#1091;&#1089; &#1080;&#1085;&#1076;&#1077;&#1082;&#1089;&#1072;&#1094;&#1080;&#1080;</div><div id="doc-status-summary" class="doc-status-summary">&#1055;&#1088;&#1086;&#1074;&#1077;&#1088;&#1103;&#1102;...</div></div><div id="doc-status-list" class="doc-status-list"></div></section>
    <section class="chat-panel"><div class="panel-head"><div class="panel-title-row"><div><h2>&#1063;&#1072;&#1090; &#1087;&#1086; &#1073;&#1072;&#1079;&#1077; &#1079;&#1085;&#1072;&#1085;&#1080;&#1081;</h2><div class="small">&#1045;&#1089;&#1083;&#1080; &#1088;&#1077;&#1083;&#1077;&#1074;&#1072;&#1085;&#1090;&#1085;&#1099;&#1093; &#1080;&#1089;&#1090;&#1086;&#1095;&#1085;&#1080;&#1082;&#1086;&#1074; &#1085;&#1077;&#1090;, &#1072;&#1089;&#1089;&#1080;&#1089;&#1090;&#1077;&#1085;&#1090; &#1089;&#1082;&#1072;&#1078;&#1077;&#1090;, &#1095;&#1090;&#1086; &#1076;&#1072;&#1085;&#1085;&#1099;&#1093; &#1085;&#1077;&#1076;&#1086;&#1089;&#1090;&#1072;&#1090;&#1086;&#1095;&#1085;&#1086;.</div></div><button id="reset-chat" class="secondary" type="button">&#1057;&#1073;&#1088;&#1086;&#1089;&#1080;&#1090;&#1100;</button></div></div><div id="messages"><div class="message"><span class="role">&#1057;&#1080;&#1089;&#1090;&#1077;&#1084;&#1072;</span>AI &#1087;&#1086;&#1080;&#1089;&#1082; &#1075;&#1086;&#1090;&#1086;&#1074;. &#1052;&#1086;&#1078;&#1085;&#1086; &#1079;&#1072;&#1076;&#1072;&#1090;&#1100; &#1074;&#1086;&#1087;&#1088;&#1086;&#1089; &#1087;&#1086; &#1073;&#1072;&#1079;&#1077; &#1079;&#1085;&#1072;&#1085;&#1080;&#1081;.</div></div></section>
    <form class="chat" id="chat"><textarea id="prompt" placeholder="&#1057;&#1087;&#1088;&#1086;&#1089;&#1080; &#1087;&#1086; &#1073;&#1072;&#1079;&#1077; &#1079;&#1085;&#1072;&#1085;&#1080;&#1081;"></textarea><button id="send" type="submit">&#1054;&#1090;&#1087;&#1088;&#1072;&#1074;&#1080;&#1090;&#1100;</button></form>
  </main>
  <script>
    const basePath="/ai";const messages=document.querySelector("#messages");const statusEl=document.querySelector("#status");const promptEl=document.querySelector("#prompt");const send=document.querySelector("#send");const resetChat=document.querySelector("#reset-chat");const docStatusSummary=document.querySelector("#doc-status-summary");const docStatusList=document.querySelector("#doc-status-list");const chatForm=document.querySelector("#chat");const storageKey="outline-ai-chat-history-v1";let chatHistory=[];
    function setStatus(text){if(statusEl)statusEl.textContent=text}
    function saveHistory(){try{localStorage.setItem(storageKey,JSON.stringify(chatHistory.slice(-80)))}catch{}}
    function loadHistory(){try{const saved=JSON.parse(localStorage.getItem(storageKey)||"[]");if(Array.isArray(saved))chatHistory=saved.filter((item)=>item&&typeof item.role==="string"&&typeof item.text==="string")}catch{chatHistory=[]}}
    async function requestJson(url,options={}){const response=await fetch(basePath+url,options);if(!response.ok){let detail=await response.text();try{detail=JSON.parse(detail).detail||detail}catch{}throw new Error(detail)}return response.json()}
    function renderMessages(){messages.innerHTML="";if(!chatHistory.length){chatHistory.push({role:"\u0421\u0438\u0441\u0442\u0435\u043c\u0430",text:"AI \u043f\u043e\u0438\u0441\u043a \u0433\u043e\u0442\u043e\u0432. \u041c\u043e\u0436\u043d\u043e \u0437\u0430\u0434\u0430\u0442\u044c \u0432\u043e\u043f\u0440\u043e\u0441 \u043f\u043e \u0431\u0430\u0437\u0435 \u0437\u043d\u0430\u043d\u0438\u0439.",sources:[]})}chatHistory.forEach((message)=>renderMessage(message.role,message.text,message.sources||[]));messages.scrollTop=messages.scrollHeight}
    function renderMessage(role,text,sources=[]){const item=document.createElement("div");item.className="message";const label=document.createElement("span");label.className="role";label.textContent=role;item.append(label,document.createTextNode(text||""));if(sources.length){const list=document.createElement("div");list.className="sources";sources.forEach((source)=>{if(!source.url)return;const link=document.createElement("a");link.href=source.url;link.target="_top";const kind=source.sourceType==="outline_attachment"?"\u0421\u0442\u0440\u0430\u043d\u0438\u0446\u0430":(source.sourceType==="file"?"\u0424\u0430\u0439\u043b":"Outline");link.textContent=`${kind}: ${source.title||source.url}`;list.append(link)});item.append(list)}messages.append(item)}
    function addMessage(role,text,sources=[]){chatHistory.push({role,text:text||"",sources:Array.isArray(sources)?sources:[]});renderMessages();saveHistory()}
    function statusLabel(status){return {indexed:"\u0418\u043d\u0434\u0435\u043a\u0441\u0438\u0440\u043e\u0432\u0430\u043d",done:"\u0418\u043d\u0434\u0435\u043a\u0441\u0438\u0440\u043e\u0432\u0430\u043d",pending:"\u041e\u0436\u0438\u0434\u0430\u0435\u0442",indexing:"\u0418\u043d\u0434\u0435\u043a\u0441\u0438\u0440\u0443\u0435\u0442\u0441\u044f",running:"\u0418\u043d\u0434\u0435\u043a\u0441\u0438\u0440\u0443\u0435\u0442\u0441\u044f",error:"\u041e\u0448\u0438\u0431\u043a\u0430",needs_ocr:"\u041d\u0443\u0436\u0435\u043d OCR",skipped:"\u041d\u0435\u0442 \u0442\u0435\u043a\u0441\u0442\u0430"}[status]||status}
    function renderDocumentStatus(data){const summary=data.summary||{};docStatusSummary.textContent=`${summary.indexed||0}/${summary.total||0} \u0438\u043d\u0434\u0435\u043a\u0441\u0438\u0440\u043e\u0432\u0430\u043d\u043e${summary.pending?`, ${summary.pending} \u043e\u0436\u0438\u0434\u0430\u0435\u0442`:""}${summary.indexing?`, ${summary.indexing} \u0432 \u0440\u0430\u0431\u043e\u0442\u0435`:""}${summary.error?`, ${summary.error} \u043e\u0448\u0438\u0431\u043e\u043a`:""}${summary.needs_ocr?`, ${summary.needs_ocr} OCR`:""}`;docStatusList.innerHTML="";(data.items||[]).slice(0,20).forEach((doc)=>{const row=document.createElement("div");row.className="doc-row";const name=document.createElement("div");name.className="doc-name";name.title=[doc.title,doc.collection,doc.error].filter(Boolean).join(" | ");name.textContent=`${doc.type==="file"?"\u0424\u0430\u0439\u043b":"\u0421\u0442\u0440\u0430\u043d\u0438\u0446\u0430"}: ${doc.title||"\u0411\u0435\u0437 \u043d\u0430\u0437\u0432\u0430\u043d\u0438\u044f"}`;const badge=document.createElement("span");badge.className=`badge ${doc.status||"pending"}`;badge.textContent=statusLabel(doc.status||"pending");row.append(name,badge);docStatusList.append(row)});if(!docStatusList.children.length){docStatusList.innerHTML='<div class="small">\u0414\u043e\u043a\u0443\u043c\u0435\u043d\u0442\u044b \u043f\u043e\u043a\u0430 \u043d\u0435 \u043d\u0430\u0439\u0434\u0435\u043d\u044b.</div>'}}
    async function refreshDocumentStatus(){try{renderDocumentStatus(await requestJson("/document-status"))}catch(error){docStatusSummary.textContent=`\u041e\u0448\u0438\u0431\u043a\u0430 \u0441\u0442\u0430\u0442\u0443\u0441\u0430: ${error.message}`}}
    async function refreshStatus(){await refreshDocumentStatus()}
    resetChat.addEventListener("click",()=>{chatHistory=[];try{localStorage.removeItem(storageKey)}catch{}renderMessages();promptEl.focus()});
    function requestHistory(){return chatHistory.filter((item)=>item.role==="\u0412\u044b"||item.role==="\u0410\u0441\u0441\u0438\u0441\u0442\u0435\u043d\u0442").slice(-10).map((item)=>({role:item.role==="\u0412\u044b"?"user":"assistant",content:item.text||""}))}
    async function submitMessage(){const message=promptEl.value.trim();if(!message)return;const history=requestHistory();promptEl.value="";send.disabled=true;addMessage("\u0412\u044b",message);setStatus("\u0414\u0443\u043c\u0430\u044e...");try{const result=await requestJson("/chat",{method:"POST",headers:{"Content-Type":"application/json"},body:JSON.stringify({message,limit:5,history})});addMessage("\u0410\u0441\u0441\u0438\u0441\u0442\u0435\u043d\u0442",result.answer||"\u041d\u0435\u0442 \u043e\u0442\u0432\u0435\u0442\u0430.",result.sources||[]);setStatus("\u0413\u043e\u0442\u043e\u0432\u043e.")}catch(error){addMessage("\u041e\u0448\u0438\u0431\u043a\u0430",error.message||"\u0417\u0430\u043f\u0440\u043e\u0441 \u043d\u0435 \u0432\u044b\u043f\u043e\u043b\u043d\u0435\u043d.");setStatus("\u0417\u0430\u043f\u0440\u043e\u0441 \u043d\u0435 \u0432\u044b\u043f\u043e\u043b\u043d\u0435\u043d.")}finally{send.disabled=false;promptEl.focus();refreshStatus().catch(()=>{})}}

    chatForm.addEventListener("submit",(event)=>{event.preventDefault();submitMessage()});
    promptEl.addEventListener("keydown",(event)=>{if(event.ctrlKey&&event.key==="Enter"){event.preventDefault();submitMessage()}});
    loadHistory();renderMessages();refreshStatus().catch(()=>{});setInterval(()=>refreshStatus().catch(()=>{}),15000);setInterval(()=>refreshStatus().catch(()=>{}),15000);
  </script>
</body>
</html>
"""
